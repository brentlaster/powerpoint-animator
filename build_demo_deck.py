"""Build a demo .pptx with shapes, connectors, block arrows, and text boxes.
Some slides are tagged with [auto-animate] in speaker notes so the tool
knows to process them; one slide is deliberately left unmarked.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def add_title(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    tf = tb.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.runs[0].font.size = Pt(28)
    p.runs[0].font.bold = True
    return tb


def add_box(slide, text, left, top, width, height, fill=RGBColor(0x4F, 0x81, 0xBD)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    tf = s.text_frame
    tf.text = text
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(16)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return s


def add_diamond(slide, text, left, top, width, height):
    s = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0xE6, 0x7E, 0x22)
    s.line.color.rgb = RGBColor(0xB3, 0x5E, 0x1A)
    tf = s.text_frame
    tf.text = text
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(14)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return s


def add_arrow_connector(slide, src_shape, dst_shape, side_src="right", side_dst="left",
                        head_at="end"):
    """Add a straight connector from src to dst.

    head_at="end"   → arrowhead on the line's end (tailEnd in XML) — default
    head_at="start" → arrowhead on the line's start (headEnd in XML) — reversed

    The physical direction of the arrow is always from src -> dst, but
    whether that's represented in XML as a default or reversed arrow is
    controlled by ``head_at``.  This lets us exercise both code paths.
    """
    def anchor(shape, side):
        left, top, w, h = shape.left, shape.top, shape.width, shape.height
        if side == "right":
            return left + w, top + h // 2
        if side == "left":
            return left, top + h // 2
        if side == "top":
            return left + w // 2, top
        if side == "bottom":
            return left + w // 2, top + h
        return left + w // 2, top + h // 2

    from lxml import etree
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

    x1, y1 = anchor(src_shape, side_src)
    x2, y2 = anchor(dst_shape, side_dst)
    if head_at == "start":
        # Draw the connector backwards so the begin-end of the path is at the dst
        x1, y1, x2, y2 = x2, y2, x1, y1
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line = conn.line
    line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    line.width = Pt(2.5)

    ln = conn.line._get_or_add_ln()
    # Clear any previous ends
    for tag in ("headEnd", "tailEnd"):
        for existing in ln.findall(f"a:{tag}", nsmap):
            ln.remove(existing)
    if head_at == "end":
        # arrowhead on line-end → tailEnd
        el = etree.SubElement(ln, f"{A}tailEnd")
    else:
        # arrowhead on line-start → headEnd
        el = etree.SubElement(ln, f"{A}headEnd")
    el.set("type", "triangle")
    el.set("w", "med")
    el.set("len", "med")
    return conn


def add_block_arrow(slide, left, top, width, height, direction="right", rotation=0):
    """Block arrow shape pointing in one cardinal direction.

    ``rotation`` is in degrees (clockwise).  A right-arrow at rotation=90
    physically points downward.  The tool should handle rotation when
    choosing the wipe direction.
    """
    mapping = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "left":  MSO_SHAPE.LEFT_ARROW,
        "up":    MSO_SHAPE.UP_ARROW,
        "down":  MSO_SHAPE.DOWN_ARROW,
    }
    s = slide.shapes.add_shape(mapping[direction], left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0x70, 0xAD, 0x47)
    s.line.color.rgb = RGBColor(0x50, 0x7E, 0x32)
    if rotation:
        s.rotation = float(rotation)
    return s


def add_text_note(slide, text, left, top, width, height, size=14):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = text
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.italic = True
            r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    return tb


def set_notes(slide, notes_text):
    slide.notes_slide.notes_text_frame.text = notes_text


def main(out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # -------------------------------------------------------------
    # Slide 1: linear process flow (marked)
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank)
    add_title(s1, "Slide 1 — Linear Pipeline")
    row_y = Inches(3.1)
    bw, bh = Inches(2.0), Inches(1.2)
    b1 = add_box(s1, "Ingest",  Inches(0.8),  row_y, bw, bh)
    b2 = add_box(s1, "Transform", Inches(4.3), row_y, bw, bh)
    b3 = add_box(s1, "Score",     Inches(7.8), row_y, bw, bh)
    b4 = add_box(s1, "Deliver",   Inches(11.0), row_y, Inches(1.8), bh, fill=RGBColor(0x2E, 0x75, 0xB6))
    add_arrow_connector(s1, b1, b2)
    add_arrow_connector(s1, b2, b3)
    add_arrow_connector(s1, b3, b4)
    add_text_note(s1, "Four-stage ETL pipeline used in production.",
                  Inches(0.8), Inches(5.2), Inches(9), Inches(0.6))
    set_notes(s1, "[auto-animate]\nExpect left-to-right reveal.")

    # -------------------------------------------------------------
    # Slide 2: decision / branch (marked)
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank)
    add_title(s2, "Slide 2 — Decision Flow")
    inp = add_box(s2, "Request", Inches(0.8), Inches(3.1), Inches(2.2), Inches(1.2))
    dec = add_diamond(s2, "Valid?", Inches(4.3), Inches(2.8), Inches(2.4), Inches(1.8))
    out_ok = add_box(s2, "Accept", Inches(8.5), Inches(1.6), Inches(2.2), Inches(1.2),
                     fill=RGBColor(0x70, 0xAD, 0x47))
    out_bad = add_box(s2, "Reject", Inches(8.5), Inches(4.6), Inches(2.2), Inches(1.2),
                      fill=RGBColor(0xC0, 0x50, 0x4D))
    add_arrow_connector(s2, inp, dec, side_src="right", side_dst="left")
    add_arrow_connector(s2, dec, out_ok, side_src="top", side_dst="left")
    add_arrow_connector(s2, dec, out_bad, side_src="bottom", side_dst="left")
    add_text_note(s2, "Yes flows up, No flows down.",
                  Inches(0.8), Inches(5.8), Inches(9), Inches(0.6))
    set_notes(s2, "[auto-animate]")

    # -------------------------------------------------------------
    # Slide 3: block arrows + text boxes (marked)
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank)
    add_title(s3, "Slide 3 — Before / After")
    before = add_box(s3, "Before", Inches(1.0), Inches(2.8), Inches(2.6), Inches(1.8),
                     fill=RGBColor(0xBF, 0xBF, 0xBF))
    after  = add_box(s3, "After",  Inches(9.7), Inches(2.8), Inches(2.6), Inches(1.8),
                     fill=RGBColor(0x70, 0xAD, 0x47))
    add_block_arrow(s3, Inches(4.2), Inches(3.1), Inches(5.0), Inches(1.2), direction="right")
    add_text_note(s3, "Legacy manual workflow, high error rate.",
                  Inches(0.8), Inches(4.9), Inches(3.0), Inches(1.0))
    add_text_note(s3, "Automated pipeline, 10× faster.",
                  Inches(9.5), Inches(4.9), Inches(3.0), Inches(1.0))
    set_notes(s3, "[auto-animate]")

    # -------------------------------------------------------------
    # Slide 4: UNMARKED — must be left alone
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank)
    add_title(s4, "Slide 4 — Unmarked, leave alone")
    add_box(s4, "Untouched", Inches(5.0), Inches(3.0), Inches(3.0), Inches(1.5))
    set_notes(s4, "no marker — tool should skip this slide")

    # -------------------------------------------------------------
    # Slide 5: arrowhead-end edge cases
    #   (a) a left-pointing block arrow between two shapes
    #   (b) a reversed connector (arrowhead on the line's start end)
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank)
    add_title(s5, "Slide 5 — Reversed / Left-pointing arrows")
    # Row 1: left-pointing block arrow (dst on left, src on right)
    dst_left = add_box(s5, "Received", Inches(0.8), Inches(1.8), Inches(2.4), Inches(1.2),
                       fill=RGBColor(0x4F, 0x81, 0xBD))
    add_block_arrow(s5, Inches(3.8), Inches(2.0), Inches(3.0), Inches(0.8), direction="left")
    src_right = add_box(s5, "Sent", Inches(7.5), Inches(1.8), Inches(2.4), Inches(1.2),
                        fill=RGBColor(0x4F, 0x81, 0xBD))
    # Row 2: reversed connector — arrowhead on START end of line. Physical
    # arrow points right→left even though bbox is drawn in the "normal" way.
    rev_dst = add_box(s5, "Target", Inches(0.8), Inches(4.3), Inches(2.4), Inches(1.2),
                      fill=RGBColor(0x70, 0xAD, 0x47))
    rev_src = add_box(s5, "Origin", Inches(7.5), Inches(4.3), Inches(2.4), Inches(1.2),
                      fill=RGBColor(0x70, 0xAD, 0x47))
    add_arrow_connector(s5, rev_src, rev_dst,
                        side_src="left", side_dst="right",
                        head_at="start")
    set_notes(s5, "[auto-animate]\nBoth rows reveal right-to-left")

    # -------------------------------------------------------------
    # Slide 6: fan-out from one source to three destinations
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank)
    add_title(s6, "Slide 6 — Fan-out")
    src = add_box(s6, "Source", Inches(0.8), Inches(3.1), Inches(2.2), Inches(1.2))
    d1 = add_box(s6, "Top",    Inches(8.0), Inches(1.3), Inches(2.2), Inches(1.0),
                 fill=RGBColor(0x70, 0xAD, 0x47))
    d2 = add_box(s6, "Middle", Inches(8.0), Inches(3.2), Inches(2.2), Inches(1.0),
                 fill=RGBColor(0xE6, 0x7E, 0x22))
    d3 = add_box(s6, "Bottom", Inches(8.0), Inches(5.1), Inches(2.2), Inches(1.0),
                 fill=RGBColor(0xC0, 0x50, 0x4D))
    add_arrow_connector(s6, src, d1, side_src="right", side_dst="left")
    add_arrow_connector(s6, src, d2, side_src="right", side_dst="left")
    add_arrow_connector(s6, src, d3, side_src="right", side_dst="left")
    set_notes(s6, "[auto-animate]\nFan-out: source, then Top, Middle, Bottom in order")

    prs.save(out_path)
    print(f"wrote {out_path}")


if __name__ == "__main__":
    import sys
    out = sys.argv[1] if len(sys.argv) > 1 else "demo_deck.pptx"
    main(out)
