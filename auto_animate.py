"""auto_animate.py — automatically add entrance animations to PowerPoint slides.

Reads a .pptx, looks at the shapes on each slide marked with "[auto-animate]"
in the speaker notes, infers a sensible click-driven animation order using
spatial flow (arrows connect shape A to shape B → reveal A, arrow, then B),
and writes a new .pptx with entrance animations injected as <p:timing> XML.

Rules of thumb:
  - Regular shapes     → Fade entrance
  - Text boxes         → Wipe from top (top-down reveal)
  - Arrows / connectors→ Wipe in the direction of the arrow
  - Flow ordering      : if arrow goes A → B, reveal A, then arrow, then B
  - Un-connected items : fall back to reading order (top-to-bottom, left-to-right)
  - Slides that already have a <p:timing> block are left alone
  - Slides whose notes do NOT contain the marker are left alone

Usage
-----
    python3 auto_animate.py input.pptx [output.pptx] [--marker TEXT] [--all]

  --marker TEXT   change the marker string (default: "[auto-animate]")
  --all           ignore markers and process every slide

"""

from __future__ import annotations

import argparse
import copy
import itertools
import os
import sys
from dataclasses import dataclass, field
from typing import Iterable, Optional

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ---------------------------------------------------------------------------
# XML namespaces
# ---------------------------------------------------------------------------
NS = {
    "a":  "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p":  "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r":  "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}
A  = "{%s}" % NS["a"]
P  = "{%s}" % NS["p"]


# ---------------------------------------------------------------------------
# Shape classification & geometry
# ---------------------------------------------------------------------------
@dataclass
class ShapeInfo:
    sp_id: int                  # PowerPoint drawingML id
    name: str
    kind: str                   # 'text', 'shape', 'connector', 'block_arrow'
    left: int
    top: int
    width: int
    height: int
    # For connectors / block arrows: direction vector (dx, dy) where the
    # HEAD of the arrow points.  Normalised to unit-ish ints (dominant axis).
    direction: Optional[tuple[int, int]] = None
    # Endpoints in slide coordinates (EMU): (x1,y1) = tail, (x2,y2) = head.
    tail: Optional[tuple[int, int]] = None
    head: Optional[tuple[int, int]] = None

    @property
    def center(self) -> tuple[int, int]:
        return (self.left + self.width // 2, self.top + self.height // 2)


# Block-arrow auto-shape geometries we know.
# These come from the `prst` attribute on <a:prstGeom> inside <p:sp>.
BLOCK_ARROW_PRST = {
    "rightArrow":  (1,  0),
    "leftArrow":   (-1, 0),
    "upArrow":     (0, -1),
    "downArrow":   (0,  1),
    "bentUpArrow": (0, -1),
    "curvedRightArrow": (1, 0),
    "curvedLeftArrow":  (-1, 0),
    "curvedUpArrow":    (0, -1),
    "curvedDownArrow":  (0,  1),
    "stripedRightArrow": (1, 0),
    "notchedRightArrow": (1, 0),
}


def _preset_geom(sp: etree._Element) -> Optional[str]:
    el = sp.find(f".//{A}prstGeom")
    return el.get("prst") if el is not None else None


def _end_arrow_kind(sp: etree._Element, end_name: str) -> str:
    """Return the arrowhead type at one end of a line/connector.

    end_name is "tailEnd" (line's end point) or "headEnd" (line's start point).
    Returns "none" if there is no element, or the element's type attribute
    (default "none" per OOXML).
    """
    el = sp.find(f".//{A}ln/{A}{end_name}")
    if el is None:
        return "none"
    return el.get("type", "none")


def _has_arrow_at(sp: etree._Element, end_name: str) -> bool:
    return _end_arrow_kind(sp, end_name) != "none"


def _xfrm(sp: etree._Element) -> Optional[etree._Element]:
    # connector shapes use p:cxnSp/p:spPr/a:xfrm ; normal shapes use p:sp/p:spPr/a:xfrm
    return sp.find(f".//{P}spPr/{A}xfrm")


def classify_shape(sp: etree._Element) -> Optional[ShapeInfo]:
    """Return ShapeInfo for a <p:sp> / <p:cxnSp>, or None if it should be skipped."""
    tag = etree.QName(sp.tag).localname
    if tag not in ("sp", "cxnSp", "pic"):
        return None

    nvSpPr = sp.find(f".//{P}nvSpPr/{P}cNvPr")
    if nvSpPr is None:
        nvSpPr = sp.find(f".//{P}nvCxnSpPr/{P}cNvPr")
    if nvSpPr is None:
        nvSpPr = sp.find(f".//{P}nvPicPr/{P}cNvPr")
    if nvSpPr is None:
        return None
    sp_id = int(nvSpPr.get("id"))
    name = nvSpPr.get("name", "")

    xfrm = _xfrm(sp)
    if xfrm is None:
        return None
    off = xfrm.find(f"{A}off")
    ext = xfrm.find(f"{A}ext")
    if off is None or ext is None:
        return None
    left = int(off.get("x"))
    top = int(off.get("y"))
    width = int(ext.get("cx"))
    height = int(ext.get("cy"))
    flipH = xfrm.get("flipH") == "1"
    flipV = xfrm.get("flipV") == "1"

    prst = _preset_geom(sp)

    # Connector (line)
    if tag == "cxnSp":
        # Endpoints inferred from bbox + flip flags
        x1, y1 = (left + width, top) if flipH else (left, top)
        x2, y2 = (left, top + height) if flipH else (left + width, top + height)
        if flipV:
            y1, y2 = y2, y1

        # Figure out which end of the line carries the arrowhead.
        # In OOXML a line has two decorations:
        #   <a:headEnd> — decoration at the LINE START  (x1,y1 as computed above)
        #   <a:tailEnd> — decoration at the LINE END    (x2,y2 as computed above)
        # An end counts as "an arrowhead" iff the element is present AND its
        # type attribute is something other than "none" (missing element or
        # type="none" both mean flat end).
        head_arrow = _has_arrow_at(sp, "headEnd")  # arrowhead at (x1,y1)?
        tail_arrow = _has_arrow_at(sp, "tailEnd")  # arrowhead at (x2,y2)?
        # Canonicalise: make (x2,y2) the pointy end.
        #   tail_arrow only             → already correct
        #   head_arrow only             → reverse
        #   both or neither             → fall back to bbox direction
        if head_arrow and not tail_arrow:
            x1, y1, x2, y2 = x2, y2, x1, y1

        dx = x2 - x1
        dy = y2 - y1
        # Normalise to dominant axis
        if abs(dx) >= abs(dy):
            direction = (1 if dx >= 0 else -1, 0)
        else:
            direction = (0, 1 if dy >= 0 else -1)

        return ShapeInfo(
            sp_id=sp_id, name=name, kind="connector",
            left=left, top=top, width=width, height=height,
            direction=direction, tail=(x1, y1), head=(x2, y2),
        )

    # Block-arrow auto-shape
    if prst in BLOCK_ARROW_PRST:
        dx, dy = BLOCK_ARROW_PRST[prst]
        if flipH:
            dx = -dx
        if flipV:
            dy = -dy
        # Apply rotation: xfrm @rot is in 60000ths of a degree, clockwise.
        # We snap to the nearest 90° since our direction vector is axis-aligned.
        rot_attr = xfrm.get("rot")
        if rot_attr:
            deg = (int(rot_attr) / 60000.0) % 360.0
            quarter = int(round(deg / 90.0)) % 4
            for _ in range(quarter):
                # rotate (dx, dy) 90° clockwise: (dx, dy) -> (-dy, dx)
                dx, dy = -dy, dx
        cx, cy = left + width // 2, top + height // 2
        # Tail and head along the direction axis (after rotation)
        if abs(dx) >= abs(dy):
            tail = (left if dx > 0 else left + width, cy)
            head = (left + width if dx > 0 else left, cy)
        else:
            tail = (cx, top if dy > 0 else top + height)
            head = (cx, top + height if dy > 0 else top)
        return ShapeInfo(
            sp_id=sp_id, name=name, kind="block_arrow",
            left=left, top=top, width=width, height=height,
            direction=(dx, dy), tail=tail, head=head,
        )

    # Text box
    if tag == "sp" and sp.find(f"{P}nvSpPr/{P}cNvSpPr[@txBox='1']") is not None:
        return ShapeInfo(sp_id=sp_id, name=name, kind="text",
                         left=left, top=top, width=width, height=height)

    # Picture
    if tag == "pic":
        return ShapeInfo(sp_id=sp_id, name=name, kind="shape",
                         left=left, top=top, width=width, height=height)

    # Default: regular shape
    return ShapeInfo(sp_id=sp_id, name=name, kind="shape",
                     left=left, top=top, width=width, height=height)


# ---------------------------------------------------------------------------
# Spatial flow inference
# ---------------------------------------------------------------------------
def _point_in_rect(pt, rect, pad=0):
    x, y = pt
    l, t, w, h = rect
    return (l - pad) <= x <= (l + w + pad) and (t - pad) <= y <= (t + h + pad)


def _dist2(a, b):
    return (a[0] - b[0]) ** 2 + (a[1] - b[1]) ** 2


def find_arrow_target(arrow: ShapeInfo,
                      candidates: list[ShapeInfo],
                      which: str) -> Optional[ShapeInfo]:
    """Return the candidate shape that best matches the arrow's tail or head."""
    pt = arrow.tail if which == "tail" else arrow.head
    if pt is None:
        return None
    # First, any candidate whose bbox contains the point (with small padding)
    pad = 50_000  # ~0.05 inch in EMU
    containing = [c for c in candidates if _point_in_rect(pt, (c.left, c.top, c.width, c.height), pad)]
    if containing:
        # prefer the smallest containing rect (most specific)
        containing.sort(key=lambda c: c.width * c.height)
        return containing[0]
    # Otherwise nearest centre within a reasonable distance (~2.5 inch)
    if not candidates:
        return None
    nearest = min(candidates, key=lambda c: _dist2(pt, c.center))
    max_d2 = (2_200_000) ** 2  # a bit over 2.3 inches in EMU
    return nearest if _dist2(pt, nearest.center) < max_d2 else None


def build_flow(shapes: list[ShapeInfo]) -> list[list[ShapeInfo]]:
    """Return a list of animation groups (each group fires on its own click).
    Each group is a list of ShapeInfo in the order they should be revealed
    within that click.
    """
    arrows = [s for s in shapes if s.kind in ("connector", "block_arrow")]
    non_arrows = [s for s in shapes if s.kind not in ("connector", "block_arrow")]
    arrow_by_id = {a.sp_id: a for a in arrows}

    # Build source/dest for every arrow
    edges: dict[int, tuple[Optional[ShapeInfo], Optional[ShapeInfo]]] = {}
    for a in arrows:
        src = find_arrow_target(a, non_arrows, "tail")
        dst = find_arrow_target(a, non_arrows, "head")
        edges[a.sp_id] = (src, dst)

    outgoing: dict[int, list[tuple[ShapeInfo, ShapeInfo]]] = {}
    incoming_count: dict[int, int] = {s.sp_id: 0 for s in non_arrows}
    for a_id, (src, dst) in edges.items():
        if src and dst:
            outgoing.setdefault(src.sp_id, []).append((arrow_by_id[a_id], dst))
            incoming_count[dst.sp_id] = incoming_count.get(dst.sp_id, 0) + 1

    # Within a single source's fan-out, visit destinations in reading order
    # (top-to-bottom, then left-to-right).  This matters when one shape has
    # multiple arrows leading to several siblings.
    for src_id, edge_list in outgoing.items():
        edge_list.sort(key=lambda ad: (ad[1].top, ad[1].left))

    # A shape is "in-flow" if it has any incoming or outgoing flow edges
    in_flow_ids = set(outgoing.keys())
    for c in incoming_count:
        if incoming_count[c] > 0:
            in_flow_ids.add(c)

    visited_arrows: set[int] = set()
    visited_shapes: set[int] = set()
    processed_outgoing: set[int] = set()
    flow_groups: list[list[ShapeInfo]] = []

    reading_order = sorted(non_arrows, key=lambda s: (s.top, s.left))

    def _bfs(root: ShapeInfo):
        queue = [root]
        while queue:
            node = queue.pop(0)
            if node.sp_id not in visited_shapes:
                visited_shapes.add(node.sp_id)
                flow_groups.append([node])
            if node.sp_id in processed_outgoing:
                continue
            processed_outgoing.add(node.sp_id)
            for arrow, dst in outgoing.get(node.sp_id, []):
                if arrow.sp_id in visited_arrows:
                    continue
                visited_arrows.add(arrow.sp_id)
                group: list[ShapeInfo] = [arrow]
                if dst.sp_id not in visited_shapes:
                    group.append(dst)
                    visited_shapes.add(dst.sp_id)
                queue.append(dst)
                flow_groups.append(group)

    # Phase 1: BFS from true sources (incoming=0, outgoing>0) in reading order
    sources = [s for s in reading_order
               if s.sp_id in outgoing and incoming_count.get(s.sp_id, 0) == 0]
    for src in sources:
        if src.sp_id not in visited_shapes:
            _bfs(src)

    # Phase 2: any unreached in-flow node (e.g. cycles) — start BFS from there
    for s in reading_order:
        if s.sp_id in in_flow_ids and s.sp_id not in visited_shapes:
            _bfs(s)

    # Phase 3: leftover dangling arrows (source or target not detected)
    leftover_arrows = [a for a in sorted(arrows, key=lambda s: (s.top, s.left))
                       if a.sp_id not in visited_arrows]
    for a in leftover_arrows:
        flow_groups.append([a])
        visited_arrows.add(a.sp_id)

    # Phase 4: orphans — text boxes / shapes with no flow edges
    orphans = [s for s in reading_order if s.sp_id not in visited_shapes]
    # Split orphans into "top-of-slide" (go first) and "below-flow" (go last)
    # using the topmost flow element as the divider.
    if flow_groups:
        flow_tops = [sh.top for grp in flow_groups for sh in grp
                     if sh.kind not in ("connector", "block_arrow")]
        divider = min(flow_tops) if flow_tops else 0
    else:
        divider = 0
    pre_orphans  = [s for s in orphans if s.top <  divider]
    post_orphans = [s for s in orphans if s.top >= divider]

    click_groups: list[list[ShapeInfo]] = []
    for s in pre_orphans:
        click_groups.append([s])
        visited_shapes.add(s.sp_id)
    click_groups.extend(flow_groups)
    for s in post_orphans:
        click_groups.append([s])
        visited_shapes.add(s.sp_id)

    return click_groups


# ---------------------------------------------------------------------------
# Animation choice
# ---------------------------------------------------------------------------
# presetID values (PresentationML entrance effects)
PRESET_APPEAR = 1    # Appear
PRESET_FADE   = 10   # Fade
PRESET_WIPE   = 12   # Wipe

# Wipe subtypes/filter values keyed by MOTION DIRECTION — the direction
# the wipe travels.  For an arrow pointing right we want motion="right":
# the reveal starts at the LEFT edge (the arrow's tail) and sweeps to the
# right (toward the arrowhead).
#
# Per ECMA-376, Wipe (presetID=12) subtypes describe where the wipe
# STARTS FROM:
#     subtype 1 = from bottom  → motion up
#     subtype 2 = from left    → motion right
#     subtype 4 = from top     → motion down
#     subtype 8 = from right   → motion left
# The <p:animEffect filter="wipe(X)"> value uses the same motion-direction
# convention as the dict key here.
WIPE_SUBTYPE    = {"right": 2, "left": 8, "down": 4, "up": 1}
WIPE_FILTER_DIR = {"right": "right", "left": "left", "down": "down", "up": "up"}


def pick_animation(s: ShapeInfo) -> dict:
    """Return a dict describing the entrance effect for this shape."""
    if s.kind in ("connector", "block_arrow"):
        # Arrow direction vector -> reveal motion follows the arrow.
        dx, dy = s.direction or (1, 0)
        if abs(dx) >= abs(dy):
            motion = "right" if dx > 0 else "left"
        else:
            motion = "down"  if dy > 0 else "up"
        return {"kind": "wipe", "motion": motion, "duration": 700}
    if s.kind == "text":
        # Text boxes reveal top-to-bottom (motion = down)
        return {"kind": "wipe", "motion": "down", "duration": 500}
    # Regular shapes & pictures
    return {"kind": "fade", "duration": 400}


# ---------------------------------------------------------------------------
# Timing-XML generation
# ---------------------------------------------------------------------------
class _IdGen:
    def __init__(self, start=1):
        self._n = itertools.count(start)

    def next(self) -> int:
        return next(self._n)


def _sub(parent, tag, **attrs):
    el = etree.SubElement(parent, tag)
    for k, v in attrs.items():
        el.set(k, str(v))
    return el


def _build_effect(effect: dict, sp_id: int, ids: _IdGen,
                  node_type: str = "withEffect") -> etree._Element:
    """Build a <p:par> node for a single entrance effect."""
    kind = effect["kind"]
    if kind == "appear":
        preset_id, preset_sub, preset_class = PRESET_APPEAR, 0, "entr"
    elif kind == "fade":
        preset_id, preset_sub, preset_class = PRESET_FADE, 0, "entr"
    elif kind == "wipe":
        preset_id, preset_sub, preset_class = PRESET_WIPE, WIPE_SUBTYPE[effect["motion"]], "entr"
    else:
        preset_id, preset_sub, preset_class = PRESET_APPEAR, 0, "entr"

    outer_par = etree.Element(f"{P}par")
    c0 = _sub(outer_par, f"{P}cTn",
              id=ids.next(), presetID=preset_id, presetClass=preset_class,
              presetSubtype=preset_sub, fill="hold", grpId=0, nodeType=node_type)
    stc = _sub(c0, f"{P}stCondLst")
    _sub(stc, f"{P}cond", delay=0)
    child_tn = _sub(c0, f"{P}childTnLst")

    # 1. Always flip visibility on via <p:set>
    set_el = _sub(child_tn, f"{P}set")
    cbhvr = _sub(set_el, f"{P}cBhvr")
    ct = _sub(cbhvr, f"{P}cTn", id=ids.next(), dur=1, fill="hold")
    st = _sub(ct, f"{P}stCondLst")
    _sub(st, f"{P}cond", delay=0)
    tgt = _sub(cbhvr, f"{P}tgtEl")
    _sub(tgt, f"{P}spTgt", spid=sp_id)
    attr_lst = _sub(cbhvr, f"{P}attrNameLst")
    attr_name = etree.SubElement(attr_lst, f"{P}attrName")
    attr_name.text = "style.visibility"
    to = _sub(set_el, f"{P}to")
    _sub(to, f"{P}strVal", val="visible")

    # 2. Add the actual motion animation (fade opacity or wipe filter)
    dur = effect.get("duration", 500)
    if kind == "fade":
        anim = _sub(child_tn, f"{P}anim", calcmode="lin", valueType="num")
        cb = _sub(anim, f"{P}cBhvr", additive="base")
        _sub(cb, f"{P}cTn", id=ids.next(), dur=dur, fill="hold")
        tgt = _sub(cb, f"{P}tgtEl")
        _sub(tgt, f"{P}spTgt", spid=sp_id)
        al = _sub(cb, f"{P}attrNameLst")
        an = etree.SubElement(al, f"{P}attrName")
        an.text = "style.opacity"
        tav = _sub(anim, f"{P}tavLst")
        tav1 = _sub(tav, f"{P}tav", tm=0)
        val1 = _sub(tav1, f"{P}val")
        _sub(val1, f"{P}fltVal", val=0)
        tav2 = _sub(tav, f"{P}tav", tm=100000)
        val2 = _sub(tav2, f"{P}val")
        _sub(val2, f"{P}fltVal", val=1)
    elif kind == "wipe":
        ae = _sub(child_tn, f"{P}animEffect", transition="in",
                  filter=f"wipe({WIPE_FILTER_DIR[effect['motion']]})")
        cb = _sub(ae, f"{P}cBhvr")
        _sub(cb, f"{P}cTn", id=ids.next(), dur=dur)
        tgt = _sub(cb, f"{P}tgtEl")
        _sub(tgt, f"{P}spTgt", spid=sp_id)
    # 'appear' needs nothing beyond the <p:set> above

    return outer_par


def _wrap_click(group: list[ShapeInfo], ids: _IdGen) -> etree._Element:
    """Wrap a list of shapes into one click-triggered animation group.
    Builds outer→inner IDs so the document-order numbering matches what
    PowerPoint produces natively.
    """
    par1 = etree.Element(f"{P}par")
    c1 = _sub(par1, f"{P}cTn", id=ids.next(), fill="hold")
    st1 = _sub(c1, f"{P}stCondLst")
    _sub(st1, f"{P}cond", delay="indefinite")
    ct1 = _sub(c1, f"{P}childTnLst")

    par2 = _sub(ct1, f"{P}par")
    c2 = _sub(par2, f"{P}cTn", id=ids.next(), fill="hold")
    st2 = _sub(c2, f"{P}stCondLst")
    _sub(st2, f"{P}cond", delay=0)
    ct2 = _sub(c2, f"{P}childTnLst")

    for i, sh in enumerate(group):
        node_type = "clickEffect" if i == 0 else "withEffect"
        eff_par = _build_effect(pick_animation(sh), sh.sp_id, ids, node_type)
        ct2.append(eff_par)
    return par1


def build_timing(click_groups: list[list[ShapeInfo]]) -> tuple[etree._Element, list[int]]:
    """Build a <p:timing> element and return (element, list of animated sp_ids)."""
    ids = _IdGen(start=2)
    root_par = etree.Element(f"{P}par")
    root_cTn = _sub(root_par, f"{P}cTn", id=ids.next(),
                    dur="indefinite", restart="never", nodeType="tmRoot")
    root_children = _sub(root_cTn, f"{P}childTnLst")

    seq = _sub(root_children, f"{P}seq", concurrent=1, nextAc="seek")
    seq_cTn = _sub(seq, f"{P}cTn", id=ids.next(),
                   dur="indefinite", nodeType="mainSeq")
    seq_children = _sub(seq_cTn, f"{P}childTnLst")

    animated_ids: list[int] = []
    for group in click_groups:
        click = _wrap_click(group, ids)
        seq_children.append(click)
        for sh in group:
            animated_ids.append(sh.sp_id)

    # prev/next condition lists (standard)
    prev_cond = _sub(seq, f"{P}prevCondLst")
    pc = _sub(prev_cond, f"{P}cond", evt="onPrev", delay=0)
    ptgt = _sub(pc, f"{P}tgtEl")
    _sub(ptgt, f"{P}sldTgt")

    next_cond = _sub(seq, f"{P}nextCondLst")
    nc = _sub(next_cond, f"{P}cond", evt="onNext", delay=0)
    ntgt = _sub(nc, f"{P}tgtEl")
    _sub(ntgt, f"{P}sldTgt")

    tnLst_par = etree.Element(f"{P}tnLst")
    tnLst_par.append(root_par)

    timing = etree.Element(f"{P}timing")
    timing.append(tnLst_par)

    bld_lst = etree.SubElement(timing, f"{P}bldLst")
    for sp_id in animated_ids:
        _sub(bld_lst, f"{P}bldP", spid=sp_id, grpId=0)

    return timing, animated_ids


# ---------------------------------------------------------------------------
# Slide-level driver
# ---------------------------------------------------------------------------
def slide_shapes_xml(slide) -> list[etree._Element]:
    """Return top-level shape elements under spTree in document order."""
    spTree = slide._element.find(f".//{P}cSld/{P}spTree")
    if spTree is None:
        return []
    result = []
    for child in spTree:
        tag = etree.QName(child.tag).localname
        if tag in ("sp", "cxnSp", "pic", "graphicFrame", "grpSp"):
            result.append(child)
    return result


def slide_has_timing(slide) -> bool:
    return slide._element.find(f"{P}timing") is not None


def notes_text(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text or ""


def process_slide(slide, verbose=True) -> int:
    """Add animations to one slide. Returns number of shapes animated."""
    shapes_xml = slide_shapes_xml(slide)
    shape_infos: list[ShapeInfo] = []
    for sp in shapes_xml:
        info = classify_shape(sp)
        if info is not None:
            shape_infos.append(info)

    if not shape_infos:
        return 0

    groups = build_flow(shape_infos)
    timing, animated_ids = build_timing(groups)

    # Remove existing <p:timing> (we checked earlier it's absent, but be safe)
    sld = slide._element
    existing = sld.find(f"{P}timing")
    if existing is not None:
        sld.remove(existing)

    # Per OOXML: cSld, clrMapOvr?, transition?, timing?, extLst?
    # Insert timing before any <p:extLst> if present, else append.
    ext = sld.find(f"{P}extLst")
    if ext is not None:
        ext.addprevious(timing)
    else:
        sld.append(timing)

    if verbose:
        click_desc = []
        for g in groups:
            click_desc.append(" + ".join(f"{s.name}({s.kind})" for s in g))
        print("  " + "\n  ".join(f"click {i+1}: {d}" for i, d in enumerate(click_desc)))
    return len(animated_ids)


def process_pptx(in_path: str, out_path: str, marker: str = "[auto-animate]",
                 animate_all: bool = False, verbose: bool = True):
    prs = Presentation(in_path)
    total_animated = 0
    for i, slide in enumerate(prs.slides, 1):
        n = notes_text(slide)
        should = animate_all or (marker in n)
        if not should:
            if verbose:
                print(f"slide {i}: skipped (no marker)")
            continue
        if slide_has_timing(slide):
            if verbose:
                print(f"slide {i}: skipped (already has animations)")
            continue
        if verbose:
            print(f"slide {i}: processing")
        n_animated = process_slide(slide, verbose=verbose)
        total_animated += n_animated
    prs.save(out_path)
    if verbose:
        print(f"\nwrote {out_path}  ({total_animated} animations injected)")
    return total_animated


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------
def main():
    ap = argparse.ArgumentParser(description="Auto-add entrance animations to a .pptx.")
    ap.add_argument("input",  help="input .pptx")
    ap.add_argument("output", nargs="?",
                    help="output .pptx (default: <input>_animated.pptx)")
    ap.add_argument("--marker", default="[auto-animate]",
                    help="text in speaker notes that opts a slide in (default: [auto-animate])")
    ap.add_argument("--all", action="store_true",
                    help="ignore markers, animate every slide")
    ap.add_argument("-q", "--quiet", action="store_true")
    args = ap.parse_args()

    out = args.output
    if not out:
        base, ext = os.path.splitext(args.input)
        out = f"{base}_animated{ext}"

    process_pptx(args.input, out, marker=args.marker,
                 animate_all=args.all, verbose=not args.quiet)


if __name__ == "__main__":
    main()
